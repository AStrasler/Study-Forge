"""
Text extraction from supported study-material formats.

Supports: PDF (digital + OCR fallback), DOCX, PPTX, TXT/MD,
images (OCR), XLSX/XLS (cell text).
"""

from __future__ import annotations

from pathlib import Path
from typing import Union

from utils.logging import get_logger

logger = get_logger(__name__)

PathLike = Union[str, Path]

# If extracted PDF text is shorter than this, try OCR on rendered pages.
PDF_OCR_CHAR_THRESHOLD = 80
# Cap OCR pages so huge decks don't melt the machine.
PDF_OCR_MAX_PAGES = 40

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif"}


def extract_text(path: PathLike) -> str:
    """Dispatch to the correct extractor based on file extension."""
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(p)
    if suffix in {".docx"}:
        return _extract_docx(p)
    if suffix in {".pptx"}:
        return _extract_pptx(p)
    if suffix in {".txt", ".md", ".markdown"}:
        return p.read_text(encoding="utf-8", errors="replace")
    if suffix in IMAGE_SUFFIXES:
        return _extract_image_ocr(p)
    if suffix in {".xlsx", ".xlsm"}:
        return _extract_xlsx(p)
    if suffix == ".xls":
        return _extract_xls_legacy(p)

    raise ValueError(f"Unsupported file type: {suffix} ({p.name})")


def _extract_pdf(path: Path) -> str:
    digital = _extract_pdf_digital(path)
    if len(digital.strip()) >= PDF_OCR_CHAR_THRESHOLD:
        return digital.strip()

    logger.info(
        "PDF %s has little digital text (%d chars); trying OCR",
        path.name,
        len(digital.strip()),
    )
    ocr_text = _extract_pdf_ocr(path)
    if ocr_text.strip():
        if digital.strip():
            return (digital.strip() + "\n\n" + ocr_text.strip()).strip()
        return ocr_text.strip()
    return digital.strip()


def _extract_pdf_digital(path: Path) -> str:
    parts: list[str] = []

    # Prefer PyMuPDF when available (often better on real lecture PDFs).
    try:
        import fitz  # pymupdf

        doc = fitz.open(str(path))
        try:
            for page in doc:
                t = page.get_text("text") or ""
                if t.strip():
                    parts.append(t.strip())
        finally:
            doc.close()
        if parts:
            return "\n\n".join(parts)
    except Exception as exc:
        logger.debug("PyMuPDF digital extract failed for %s: %s", path.name, exc)

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            try:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text.strip())
            except Exception as exc:
                logger.warning("PDF page %d extraction failed in %s: %s", i, path.name, exc)
    except Exception as exc:
        logger.warning("pypdf failed for %s: %s", path.name, exc)

    return "\n\n".join(parts)


def _extract_pdf_ocr(path: Path) -> str:
    """Render PDF pages to images and OCR them."""
    try:
        import fitz  # pymupdf
    except ImportError as exc:
        logger.warning("OCR PDF needs pymupdf: %s", exc)
        return ""

    try:
        _require_ocr()
    except RuntimeError as exc:
        logger.warning("%s", exc)
        return ""

    parts: list[str] = []
    doc = fitz.open(str(path))
    try:
        n = min(len(doc), PDF_OCR_MAX_PAGES)
        for i in range(n):
            page = doc[i]
            # 2x scale improves OCR without huge memory
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            png_bytes = pix.tobytes("png")
            text = _ocr_png_bytes(png_bytes)
            if text.strip():
                parts.append(text.strip())
                logger.info("OCR PDF page %d/%d (%s): %d chars", i + 1, n, path.name, len(text.strip()))
    finally:
        doc.close()

    return "\n\n".join(parts)


def _extract_image_ocr(path: Path) -> str:
    try:
        _require_ocr()
    except RuntimeError as exc:
        raise RuntimeError(
            f"Cannot OCR {path.name}: {exc}. "
            "On Deepnote run: apt-get update && apt-get install -y tesseract-ocr"
        ) from exc

    from PIL import Image

    img = Image.open(str(path))
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    text = _ocr_pil_image(img)
    if not text.strip():
        logger.warning("OCR returned empty text for %s", path.name)
    return text.strip()


def _require_ocr() -> None:
    try:
        import pytesseract  # noqa: F401
        from PIL import Image  # noqa: F401
    except ImportError as exc:
        raise RuntimeError(
            "OCR requires pillow and pytesseract (pip install pillow pytesseract)"
        ) from exc

    import pytesseract

    try:
        pytesseract.get_tesseract_version()
    except Exception as exc:
        raise RuntimeError(
            "Tesseract binary not found. "
            "Deepnote/Debian: sudo apt-get update && sudo apt-get install -y tesseract-ocr"
        ) from exc


def _ocr_pil_image(img) -> str:
    import pytesseract

    # eng default; extend later with OCR_LANG env if needed
    return pytesseract.image_to_string(img) or ""


def _ocr_png_bytes(png_bytes: bytes) -> str:
    from io import BytesIO

    from PIL import Image

    img = Image.open(BytesIO(png_bytes))
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    return _ocr_pil_image(img)


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n\n".join(parts).strip()


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("XLSX support needs openpyxl: pip install openpyxl") from exc

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb:
            parts.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts).strip()


def _extract_xls_legacy(path: Path) -> str:
    """Old .xls — optional xlrd."""
    try:
        import xlrd
    except ImportError as exc:
        raise RuntimeError(
            "Legacy .xls needs xlrd: pip install xlrd. Prefer saving as .xlsx."
        ) from exc

    book = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"## Sheet: {sheet.name}")
        for r in range(sheet.nrows):
            cells = [
                str(sheet.cell_value(r, c)).strip()
                for c in range(sheet.ncols)
                if str(sheet.cell_value(r, c)).strip()
            ]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()
